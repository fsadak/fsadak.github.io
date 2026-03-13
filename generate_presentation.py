"""
Godot Engine Eğitim Serisi - Bölüm 1: Sunum Oluşturucu
Koyu temalı, görsellerle desteklenmiş profesyonel PowerPoint sunumu.
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os
import io
import tempfile

# --- CONFIGURATION ---
BASE_DIR = r"e:\Data\Blogs_YouTube\Godot\Blog\fsadak.github.io"
ASSETS_DIR = os.path.join(BASE_DIR, "assets", "images")
OUTPUT_FILE = os.path.join(BASE_DIR, "Godot_Bolum1_Sunum.pptx")

# Colors - Dark theme with Godot accent
BG_COLOR = RGBColor(0x1E, 0x1E, 0x2E)        # Dark background
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)       # White
SUBTITLE_COLOR = RGBColor(0xB0, 0xB0, 0xCC)   # Light gray-purple
BODY_COLOR = RGBColor(0xE0, 0xE0, 0xE8)       # Light gray
ACCENT_COLOR = RGBColor(0x47, 0x8C, 0xBF)     # Godot blue
ACCENT2_COLOR = RGBColor(0x6C, 0x9E, 0x3A)    # Green accent
CODE_BG_COLOR = RGBColor(0x11, 0x11, 0x1B)    # Darker for code blocks
HIGHLIGHT_COLOR = RGBColor(0xFF, 0xCC, 0x33)   # Yellow highlight
DIVIDER_COLOR = RGBColor(0x47, 0x8C, 0xBF)    # Blue divider

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---- HELPERS ----

def load_image_as_png(image_name):
    """Load an image (including webp) and return as PNG bytes."""
    path = os.path.join(ASSETS_DIR, image_name)
    if not os.path.exists(path):
        print(f"  [UYARI] Görsel bulunamadı: {path}")
        return None
    try:
        img = Image.open(path)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf
    except Exception as e:
        print(f"  [HATA] Görsel yüklenemedi {image_name}: {e}")
        return None


def set_slide_bg(slide, color=BG_COLOR):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    return shape


def add_divider(slide, top, left=Inches(0.8), width=Inches(2.5)):
    """Add a thin colored accent divider line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DIVIDER_COLOR
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BODY_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    """Add a text box with single-style text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left=Inches(0.8), top=Inches(2.2),
                              width=Inches(11), font_size=20, color=BODY_COLOR,
                              spacing=Pt(8)):
    """Add bulleted text items to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Support bold prefix with "**text**" syntax
        if "**" in item:
            parts = item.split("**")
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.name = "Segoe UI"
                if j % 2 == 1:  # odd indices are bold
                    run.font.bold = True
                    run.font.color.rgb = ACCENT_COLOR
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Segoe UI"

        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
    return txBox


def add_section_title(slide, title, subtitle=""):
    """Standard section title layout."""
    set_slide_bg(slide)
    # Accent bar on left
    add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_COLOR)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
                 title, font_size=36, color=TITLE_COLOR, bold=True)
    add_divider(slide, Inches(1.25))

    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
                     subtitle, font_size=18, color=SUBTITLE_COLOR)


def add_image_to_slide(slide, image_name, left, top, width=None, height=None):
    """Add an image to slide, converting webp to png if needed."""
    buf = load_image_as_png(image_name)
    if buf is None:
        return None
    try:
        kwargs = {"image_file": buf, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        if not width and not height:
            kwargs["height"] = Inches(3.5)
        pic = slide.shapes.add_picture(**kwargs)
        return pic
    except Exception as e:
        print(f"  [HATA] Görsel eklenemedi {image_name}: {e}")
        return None


def add_code_block(slide, code_text, left=Inches(0.8), top=Inches(2.5),
                   width=Inches(11), height=Inches(3.5)):
    """Add a styled code block."""
    # Background rectangle
    bg = add_shape_rect(slide, left - Inches(0.1), top - Inches(0.1),
                        width + Inches(0.2), height + Inches(0.2), CODE_BG_COLOR)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xA6, 0xE2, 0x2E)  # Green code color
        p.font.name = "Consolas"
        p.space_after = Pt(2)
    return txBox


def add_table_slide(slide, headers, rows, left=Inches(0.8), top=Inches(2.2),
                    col_widths=None):
    """Add a styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    width = Inches(11.5)
    height = Inches(0.5) * num_rows

    if not col_widths:
        col_widths = [width // num_cols] * num_cols

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Style header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = TITLE_COLOR
            paragraph.font.bold = True
            paragraph.font.name = "Segoe UI"
            paragraph.alignment = PP_ALIGN.LEFT

    # Style rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = value
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x3C)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x33)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = BODY_COLOR
                paragraph.font.name = "Segoe UI"
                paragraph.alignment = PP_ALIGN.LEFT

    return table_shape


# ---- MAIN PRESENTATION BUILDER ----

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ============================================================
    # SLIDE 1: TITLE SLIDE
    # ============================================================
    print("Slayt 1: Başlık")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    # Large accent rectangle at top
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Godot Engine Eğitim Serisi", font_size=48, color=TITLE_COLOR, bold=True)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
                 "Bölüm 1: Godot'a Giriş — Motor, Kavramlar, Editör ve Felsefe",
                 font_size=28, color=ACCENT_COLOR, bold=True)

    add_divider(slide, Inches(4.0), Inches(1), Inches(4))

    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 "Sıfırdan profesyonel düzeyde oyun geliştirmeye doğru ilk adım",
                 font_size=20, color=SUBTITLE_COLOR)

    # ============================================================
    # SLIDE 2: İÇİNDEKİLER
    # ============================================================
    print("Slayt 2: İçindekiler")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "İçindekiler")

    items = [
        "1.  Godot Nedir?",
        "2.  Godot ile Neler Yapıldı?",
        "3.  Programlama Dilleri: Hangisini Seçmeliyim?",
        "4.  Temel Kavramlar: Node, Scene, Scene Tree, Signal",
        "5.  Godot Editörüne İlk Bakış",
        "6.  Godot'nun Tasarım Felsefesi",
        "7.  Kaynaklar ve Topluluk",
        "8.  Özet ve Sıradaki Adım",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), font_size=22, color=BODY_COLOR)

    # ============================================================
    # SLIDE 3: GODOT NEDİR?
    # ============================================================
    print("Slayt 3: Godot Nedir?")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Godot Nedir?")

    items = [
        "•  Hem **2D** hem de **3D** oyun geliştirmek için tasarlanmış genel amaçlı bir oyun motoru",
        "•  Masaüstü, mobil ve web platformlarında yayınlama desteği",
        "•  Konsol desteği mevcut (güçlü programlama altyapısı gerektirir)",
        "•  **Tamamen ücretsiz** ve **açık kaynak** (MIT lisansı)",
        "•  2001'de Arjantinli bir stüdyo tarafından geliştirildi",
        "•  2014'te açık kaynak olarak yayınlandı",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0))

    # ============================================================
    # SLIDE 4: GODOT İLE NELER YAPILDI? (1/2)
    # ============================================================
    print("Slayt 4: Godot ile Neler Yapıldı?")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Godot ile Neler Yapıldı?")

    items = [
        "•  **Cassette Beasts** — yaratıkları kasetlere dönüştürdüğün RPG",
        "•  **PVKK** — top-down savaş oyunu",
        "•  **Usagi Shima** — tavşan adası simülasyonu",
        "•  **Pixelorama** — açık kaynak piksel sanat çizim programı",
        "•  **RPG in a Box** — voxel RPG oluşturma aracı",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(5.5))

    # Add screenshots on right
    add_image_to_slide(slide, "introduction_cassette_beasts.webp",
                       Inches(7), Inches(1.8), height=Inches(2.5))
    add_image_to_slide(slide, "introduction_usagi_shima.webp",
                       Inches(7), Inches(4.5), height=Inches(2.5))

    # ============================================================
    # SLIDE 5: GODOT İLE NELER YAPILDI? (2/2)
    # ============================================================
    print("Slayt 5: Oyun Örnekleri - Görseller")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Godot ile Yapılmış Projeler")

    add_image_to_slide(slide, "introduction_pvkk.webp",
                       Inches(0.5), Inches(2.0), height=Inches(2.3))
    add_image_to_slide(slide, "introduction_rpg_in_a_box.webp",
                       Inches(5.0), Inches(2.0), height=Inches(2.3))

    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(4), Inches(0.5),
                 "PVKK", font_size=14, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.0), Inches(4.5), Inches(4), Inches(0.5),
                 "RPG in a Box", font_size=14, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.5),
                 "Daha fazlası: godotengine.org/showcase",
                 font_size=16, color=ACCENT_COLOR)

    # ============================================================
    # SLIDE 6: PROGRAMLAMA DİLLERİ
    # ============================================================
    print("Slayt 6: Programlama Dilleri")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Programlama Dilleri: Hangisini Seçmeliyim?")

    # Three columns
    # GDScript
    add_shape_rect(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5),
                   RGBColor(0x28, 0x28, 0x3C))
    add_text_box(slide, Inches(0.7), Inches(2.2), Inches(3.4), Inches(0.6),
                 "GDScript", font_size=24, color=ACCENT_COLOR, bold=True)
    items_gd = [
        "•  Godot'ya özgü, Python benzeri",
        "•  Öğrenmesi en kolay seçenek",
        "•  Motorla sıkı entegrasyon",
        "•  Yeni başlayanlar için önerilen",
    ]
    add_bullet_slide_content(slide, items_gd, left=Inches(0.7), top=Inches(3.0),
                              width=Inches(3.4), font_size=16)

    # C#
    add_shape_rect(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.5),
                   RGBColor(0x28, 0x28, 0x3C))
    add_text_box(slide, Inches(4.9), Inches(2.2), Inches(3.4), Inches(0.6),
                 "C#", font_size=24, color=ACCENT2_COLOR, bold=True)
    items_cs = [
        "•  Sektörde yaygın kullanım",
        "•  Büyük projelere uygun",
        "•  Unity'den geçiş kolaylığı",
        "•  Performans avantajı",
    ]
    add_bullet_slide_content(slide, items_cs, left=Inches(4.9), top=Inches(3.0),
                              width=Inches(3.4), font_size=16)

    # GDExtension
    add_shape_rect(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.5),
                   RGBColor(0x28, 0x28, 0x3C))
    add_text_box(slide, Inches(9.1), Inches(2.2), Inches(3.4), Inches(0.6),
                 "GDExtension", font_size=24, color=HIGHLIGHT_COLOR, bold=True)
    items_ext = [
        "•  C++, Rust, D, Haxe, Swift",
        "•  Yüksek performans",
        "•  Motoru derlemeye gerek yok",
        "•  SDK entegrasyonu için ideal",
    ]
    add_bullet_slide_content(slide, items_ext, left=Inches(9.1), top=Inches(3.0),
                              width=Inches(3.4), font_size=16)

    # ============================================================
    # SLIDE 7: NEDEN GDSCRIPT?
    # ============================================================
    print("Slayt 7: Neden GDScript?")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Neden GDScript ile Başlamalısınız?")

    items = [
        "•  **Godot için tasarlandı** — motorun yapısına doğal uyum sağlar",
        "•  **Sözdizimi sade ve okunabilir** — süslü parantez yok, Python benzeri",
        "•  **Öğrenme hızı yüksek** — oyun yapmaya odaklanabilirsiniz",
        "•  **Yerleşik tipler:** Vector, Color gibi oyun geliştirmeye özel tipler",
        "•  **Otomatik tamamlama** — statik dil kalitesinde kod yardımı",
        "",
        "\"Bu eğitim serisinde biz GDScript kullanacağız.\"",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0))

    # ============================================================
    # SLIDE 8: GDSCRIPT KOD ÖRNEĞİ
    # ============================================================
    print("Slayt 8: GDScript Kod Örneği")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "GDScript — Kod Örneği", "Basit bir karakter hareketi")

    code = """extends CharacterBody2D

const SPEED = 200.0

func _physics_process(delta):
    var direction = Input.get_axis("ui_left", "ui_right")
    velocity.x = direction * SPEED
    move_and_slide()"""
    add_code_block(slide, code, top=Inches(2.5), height=Inches(3.5))

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
                 "Python'a benzer, sade ve anlaşılır sözdizimi",
                 font_size=16, color=SUBTITLE_COLOR)

    # ============================================================
    # SLIDE 9: BÖLÜM AYIRICI - TEMEL KAVRAMLAR
    # ============================================================
    print("Slayt 9: Bölüm Ayırıcı - Temel Kavramlar")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Temel Kavramlar", font_size=52, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Node  •  Scene  •  Scene Tree  •  Signal",
                 font_size=28, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
                 "Godot'un çalışma mantığını anlamak için bilmeniz gereken dört yapı taşı",
                 font_size=18, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 10: BÜYÜK RESİM
    # ============================================================
    print("Slayt 10: Büyük Resim")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Büyük Resim: Oyunun Yapısı")

    items = [
        "•  Godot'da bir oyun, iç içe geçmiş **sahnelerden** oluşan bir ağaçtır",
        "•  Bu sahneler **node'lardan** meydana gelir",
        "•  Node'lar birbiriyle **sinyal sistemi** aracılığıyla konuşur",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(5.5), font_size=20)

    add_image_to_slide(slide, "key_concepts_main_menu.webp",
                       Inches(7), Inches(2.0), height=Inches(4.5))

    # ============================================================
    # SLIDE 11: SAHNE (SCENE) NEDİR?
    # ============================================================
    print("Slayt 11: Sahne (Scene)")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Sahne (Scene) Nedir?",
                      "Yeniden kullanılabilir yapı birimleri")

    items = [
        "•  Oyuncunun karakteri",
        "•  Bir silah",
        "•  Kullanıcı arayüzündeki bir menü",
        "•  Tek bir ev ya da oda",
        "•  Bir seviyenin tamamı",
        "",
        "•  Diğer motorlardaki hem **prefab** hem de **scene** kavramının rolünü üstlenir",
        "•  Sahneler **iç içe geçebilir** (nested)",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(5.5))

    add_image_to_slide(slide, "key_concepts_scene_example.webp",
                       Inches(7), Inches(2.0), height=Inches(4.5))

    # ============================================================
    # SLIDE 12: SAHNE KOMPOZİSYONU VE KALITIM
    # ============================================================
    print("Slayt 12: Sahne Kompozisyonu")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Sahne Kompozisyonu ve Kalıtım")

    items = [
        "**Kompozisyon Örneği:**",
        "•  BlinkingLight sahnesi oluştur",
        "•  BlinkingLight'ı kullanan BrokenLantern sahnesi oluştur",
        "•  BrokenLantern'larla dolu bir şehir yarat",
        "•  BlinkingLight'ın rengini değiştir → tüm şehir güncellenir!",
        "",
        "**Kalıtım:**",
        "•  Character → Magician (Character'dan kalıtım alır)",
        "•  Ana sahnedeki değişiklik alt sahnelere yansır",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6))

    add_image_to_slide(slide, "engine_design_01.png",
                       Inches(7.5), Inches(1.8), height=Inches(2.2))
    add_image_to_slide(slide, "engine_design_02.webp",
                       Inches(7.5), Inches(4.3), height=Inches(2.5))

    # ============================================================
    # SLIDE 13: NODE (DÜĞÜM) NEDİR?
    # ============================================================
    print("Slayt 13: Node (Düğüm)")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Node (Düğüm) Nedir?",
                      "Oyunun en küçük yapı taşları")

    items = [
        "Bir karakter için örnek node yapısı:",
        "",
        "•  **CharacterBody2D** — \"Player\", fizik ve hareket",
        "    ◦  **Camera2D** — oyuncuyu takip eden kamera",
        "    ◦  **Sprite2D** — karakterin görseli",
        "    ◦  **CollisionShape2D** — çarpışma alanı",
        "",
        "•  Node'lar **bileşen (component) değildir** — birbirinden bağımsız çalışır",
        "•  Sprite2D aynı zamanda Node2D, CanvasItem ve Node'dur",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6))

    add_image_to_slide(slide, "key_concepts_character_nodes.webp",
                       Inches(7.5), Inches(2.0), height=Inches(4))

    # ============================================================
    # SLIDE 14: SAHNE AĞACI (SCENE TREE)
    # ============================================================
    print("Slayt 14: Sahne Ağacı")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Sahne Ağacı (Scene Tree)",
                      "Tüm sahnelerin oluşturduğu bütünsel yapı")

    items = [
        "•  Tüm sahneler bir araya geldiğinde **sahne ağacını** oluşturur",
        "•  Teknik olarak devasa bir **node ağacıdır**",
        "•  Pratikte **sahne** bazında düşünmek daha kolay",
        "•  Sahneler somut kavramları temsil eder:",
        "   bir karakter, bir düşman, bir kapı, bir menü...",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(5.5))

    add_image_to_slide(slide, "key_concepts_scene_tree.webp",
                       Inches(7), Inches(2.0), height=Inches(4.5))

    # ============================================================
    # SLIDE 15: SİNYAL (SIGNAL)
    # ============================================================
    print("Slayt 15: Sinyal (Signal)")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Sinyal (Signal) Nedir?",
                      "Observer Pattern'ın Godot'daki karşılığı")

    items = [
        "•  Node'lar belirli bir olay gerçekleştiğinde **sinyal yayarlar (emit)**",
        "•  Node'ları kodda doğrudan bağlamak zorunda kalmadan haberdar edebilirsiniz",
        "",
        "**Yerleşik sinyal örnekleri:**",
        "•  İki nesnenin çarpışması",
        "•  Bir karakterin alana girmesi / çıkması",
        "•  Animasyonun bitmesi",
        "•  Timer'ın dolması",
        "",
        "•  Kendi **özel sinyallerinizi** de tanımlayabilirsiniz",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6))

    add_image_to_slide(slide, "key_concepts_signals.webp",
                       Inches(7.5), Inches(2.5), height=Inches(3.5))

    # ============================================================
    # SLIDE 16: SİNYAL KOD ÖRNEĞİ
    # ============================================================
    print("Slayt 16: Sinyal Kod Örneği")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Sinyal — Kod Örneği",
                      "Buton tıklamasıyla sahne değiştirme")

    code = """func _ready():
    $Button.pressed.connect(_on_button_pressed)

func _on_button_pressed():
    get_tree().change_scene_to_file("res://game.tscn")"""
    add_code_block(slide, code, top=Inches(2.5), height=Inches(2.5))

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
                 "Buton \"pressed\" sinyalini yayar → bağlı fonksiyon otomatik çalışır",
                 font_size=18, color=SUBTITLE_COLOR)

    # ============================================================
    # SLIDE 17: DÖRT KAVRAM ÖZETİ
    # ============================================================
    print("Slayt 17: Dört Kavram Özeti")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Dört Kavramın Özeti")

    headers = ["Kavram", "Ne İşe Yarar?"]
    rows = [
        ["Node (Düğüm)", "Oyunun en küçük yapı taşı. Her node'un belirli bir görevi var."],
        ["Sahne (Scene)", "Node'ların bir araya gelmesiyle oluşan yeniden kullanılabilir birim."],
        ["Sahne Ağacı (Scene Tree)", "Oyundaki tüm sahnelerin oluşturduğu bütünsel yapı."],
        ["Sinyal (Signal)", "Node'ların olay bazlı iletişim kurma sistemi."],
    ]
    add_table_slide(slide, headers, rows)

    # ============================================================
    # SLIDE 18: BÖLÜM AYIRICI - EDİTÖR
    # ============================================================
    print("Slayt 18: Bölüm Ayırıcı - Editör")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Godot Editörüne İlk Bakış", font_size=52, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Proje Yöneticisi  •  Viewport  •  Paneller  •  5 Ana Ekran",
                 font_size=24, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 19: PROJE YÖNETİCİSİ
    # ============================================================
    print("Slayt 19: Proje Yöneticisi")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Proje Yöneticisi (Project Manager)")

    items = [
        "•  Godot'yu başlattığınızda ilk çıkan pencere",
        "•  Mevcut projeleri yönetme, yeni proje oluşturma",
        "•  Dışarıdan proje içe aktarma (import)",
        "•  **Asset Library** sekmesi: topluluk projeleri ve şablonları",
        "•  Gizlilik: varsayılan olarak çevrimdışı",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(5.5))

    add_image_to_slide(slide, "editor_intro_project_manager.webp",
                       Inches(7), Inches(2.0), height=Inches(4))

    # ============================================================
    # SLIDE 20: EDİTÖR ARAYÜZÜ
    # ============================================================
    print("Slayt 20: Editör Arayüzü")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Editör Arayüzü — Genel Görünüm")

    add_image_to_slide(slide, "editor_intro_editor_empty.webp",
                       Inches(0.5), Inches(2.0), height=Inches(5))

    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                 "Üst menü  |  Viewport  |  Paneller (FileSystem, Scene, Inspector)  |  Alt panel",
                 font_size=14, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 21: VIEWPORT VE PANELLER
    # ============================================================
    print("Slayt 21: Viewport ve Paneller")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Viewport ve Paneller")

    items = [
        "**Viewport** — sahneyi görsel olarak düzenlediğiniz alan",
        "**FileSystem Paneli** — tüm proje dosyalarına erişim",
        "**Scene Paneli** — node hiyerarşisini yönetme",
        "**Inspector Paneli** — seçili node'un özelliklerini düzenleme",
        "**Alt Panel** — animasyon editörü, konsol, ses mikseri",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(5.5), font_size=18)

    # Show panels
    add_image_to_slide(slide, "editor_intro_filesystem_dock.webp",
                       Inches(7), Inches(1.8), height=Inches(2.5))
    add_image_to_slide(slide, "editor_intro_inspector_dock.webp",
                       Inches(9.5), Inches(1.8), height=Inches(2.5))
    add_image_to_slide(slide, "editor_intro_scene_dock.webp",
                       Inches(7), Inches(4.5), height=Inches(2.5))

    # ============================================================
    # SLIDE 22: BEŞ ANA EKRAN
    # ============================================================
    print("Slayt 22: Beş Ana Ekran")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Beş Ana Ekran")

    headers = ["Ekran", "Açıklama"]
    rows = [
        ["2D", "2D oyunlar ve kullanıcı arayüzleri için"],
        ["3D", "Kameralar, ışıklar ve 3D objelerle çalışma"],
        ["Script", "Kod editörü — debugger, auto-completion dahil"],
        ["Game", "Oyunu gerçek zamanlı test etme (değişiklikler kaydedilmez)"],
        ["Asset Library", "Topluluk eklentileri ve varlıkları"],
    ]
    add_table_slide(slide, headers, rows, top=Inches(2.0))

    # ============================================================
    # SLIDE 23: 2D VE 3D EKRANLARI
    # ============================================================
    print("Slayt 23: 2D ve 3D Ekranları")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "2D ve 3D Çalışma Alanları")

    add_image_to_slide(slide, "editor_intro_workspace_2d.webp",
                       Inches(0.5), Inches(2.0), height=Inches(2.3))
    add_image_to_slide(slide, "editor_intro_workspace_3d.webp",
                       Inches(0.5), Inches(4.6), height=Inches(2.3))

    add_text_box(slide, Inches(0.5), Inches(4.35), Inches(6), Inches(0.3),
                 "2D Ekranı", font_size=14, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, Inches(0.5), Inches(6.95), Inches(6), Inches(0.3),
                 "3D Ekranı", font_size=14, color=ACCENT_COLOR, bold=True)

    # Script and Game on the right
    add_image_to_slide(slide, "editor_intro_workspace_script.webp",
                       Inches(7), Inches(2.0), height=Inches(2.3))
    add_image_to_slide(slide, "editor_intro_workspace_game.webp",
                       Inches(7), Inches(4.6), height=Inches(2.3))

    add_text_box(slide, Inches(7), Inches(4.35), Inches(6), Inches(0.3),
                 "Script Ekranı", font_size=14, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, Inches(7), Inches(6.95), Inches(6), Inches(0.3),
                 "Game Ekranı", font_size=14, color=ACCENT_COLOR, bold=True)

    # ============================================================
    # SLIDE 24: YERLEŞİK SINIF REFERANSI
    # ============================================================
    print("Slayt 24: Yerleşik Sınıf Referansı")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Yerleşik Sınıf Referansı")

    items = [
        "Erişim yolları:",
        "•  **F1** tuşuna basmak",
        "•  Script ekranında **\"Search Help\"** butonu",
        "•  **Help > Search Help** menüsü",
        "•  Sınıf/fonksiyon adına **Ctrl + Click**",
        "",
        "Referans sayfası şunları içerir:",
        "•  Kalıtım hiyerarşisi",
        "•  Sınıf özeti ve kullanım senaryoları",
        "•  Özellikler, metodlar, sinyaller",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(5.5), font_size=18)

    add_image_to_slide(slide, "editor_intro_search_help.webp",
                       Inches(7), Inches(2.0), height=Inches(2.5))
    add_image_to_slide(slide, "editor_intro_help_class_animated_sprite.webp",
                       Inches(7), Inches(4.7), height=Inches(2.5))

    # ============================================================
    # SLIDE 25: BÖLÜM AYIRICI - TASARIM FELSEFESİ
    # ============================================================
    print("Slayt 25: Bölüm Ayırıcı - Tasarım Felsefesi")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Godot'nun Tasarım Felsefesi", font_size=48, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Motor neden böyle çalışıyor?",
                 font_size=24, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 26: HER ŞEY DAHİL PAKET
    # ============================================================
    print("Slayt 26: Her Şey Dahil Paket")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Her Şey Dahil Paket")

    items = [
        "Godot yerleşik olarak şunları sunar:",
        "•  **Kod editörü** (scripting workspace)",
        "•  **Animasyon editörü**",
        "•  **Tilemap editörü**",
        "•  **Shader editörü**",
        "•  **Hata ayıklayıcı** ve **profilleyici**",
        "•  Yerel ve uzak cihazlarda **hot-reload**",
        "",
        "Harici araçlar: Blender, VSCode, Visual Studio",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6))

    add_image_to_slide(slide, "introduction_editor.webp",
                       Inches(7), Inches(2.0), height=Inches(4.5))

    # ============================================================
    # SLIDE 27: AÇIK KAYNAK VE TOPLULUK
    # ============================================================
    print("Slayt 27: Açık Kaynak ve Topluluk")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Açık Kaynak ve Topluluk Odaklı")

    items = [
        "**Açık Kaynak (MIT Lisansı):**",
        "•  Tamamen ücretsiz — hiçbir kısıtlama yok",
        "•  Kaynak kodu indirilebilir, değiştirilebilir, paylaşılabilir",
        "•  Kazandığınız para tamamen sizin!",
        "",
        "**Topluluk Odaklı:**",
        "•  Binlerce katkıcı tarafından geliştirilir",
        "•  Kullanıcı ihtiyaçları geliştirme yönünü belirler",
        "•  Her büyük sürümde farklı alanlarda iyileştirmeler",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0))

    # ============================================================
    # SLIDE 28: EDİTÖR = GODOT OYUNU + AYRI MOTORLAR
    # ============================================================
    print("Slayt 28: Editör = Godot Oyunu")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Editör Bir Godot Oyunudur + Ayrı 2D/3D Motorları")

    items = [
        "**Editör = Godot Oyunu:**",
        "•  Editör, motorun kendi UI sistemiyle çalışır",
        "•  **@tool** ile GDScript kodunu editör içinde çalıştırabilirsiniz",
        "•  Import/export eklentileri, özel editörler oluşturabilirsiniz",
        "",
        "**Ayrı 2D ve 3D Motorları:**",
        "•  2D sahnelerinin temel birimi **piksel**",
        "•  3D ortamda 2D, 2D ortamda 3D render mümkün",
        "•  2D projelerde gereksiz 3D hesaplama yükü olmaz",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6.5))

    add_image_to_slide(slide, "engine_design_03.png",
                       Inches(7.5), Inches(3.5), height=Inches(3))

    # ============================================================
    # SLIDE 29: BÖLÜM AYIRICI - KAYNAKLAR
    # ============================================================
    print("Slayt 29: Bölüm Ayırıcı - Kaynaklar")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Kaynaklar ve Topluluk", font_size=52, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Godot'da yeni şeyler nasıl öğrenilir?",
                 font_size=24, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 30: KAYNAKLAR
    # ============================================================
    print("Slayt 30: Kaynaklar")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Öğrenme Kaynakları")

    headers = ["Kaynak", "Ne İşe Yarar?"]
    rows = [
        ["Resmi Kılavuz (User Manual)", "Kavramlar, özellikler ve editör kullanımı"],
        ["Yerleşik Sınıf Referansı", "API, sınıflar, metodlar, sinyaller"],
        ["Godot Forumu", "Soru sor, cevapları ara — forum.godotengine.org"],
        ["Automate The Boring Stuff", "Python ile programlama temelleri (ücretsiz)"],
        ["Topluluk Eğitimleri", "Oyun türüne özel (RPG, platform, bulmaca) eğitimler"],
    ]
    add_table_slide(slide, headers, rows, top=Inches(2.0))

    add_image_to_slide(slide, "manual_search.png",
                       Inches(8), Inches(5.0), height=Inches(2))

    # ============================================================
    # SLIDE 31: ETKİLİ SORU SORMA
    # ============================================================
    print("Slayt 31: Etkili Soru Sorma")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Etkili Soru Sormanın Sırları")

    items = [
        "1.  **Hedefinizi açıklayın** — ne yapmaya çalıştığınızı anlatın",
        "2.  **Hata mesajının tamamını paylaşın** — Debugger panelinden kopyalayın",
        "3.  **İlgili kod parçasını paylaşın** — Pastebin kullanabilirsiniz",
        "4.  **Scene Dock ekran görüntüsünü ekleyin** — bağlamı netleştirir",
        "5.  **Oyunundan video paylaşın** — OBS Studio veya ScreenToGIF",
        "6.  **Godot sürümünü belirtin** — özellikler sürümden sürüme değişir",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), font_size=20)

    # ============================================================
    # SLIDE 32: GENEL ÖZET TABLOSU
    # ============================================================
    print("Slayt 32: Genel Özet Tablosu")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Genel Özet")

    headers = ["Özellik", "Açıklama"]
    rows = [
        ["Fiyat", "Tamamen ücretsiz"],
        ["Lisans", "MIT (ticari kullanım serbestir)"],
        ["2D / 3D Desteği", "Mükemmel 2D, güçlü ve gelişen 3D"],
        ["Script Dili", "GDScript, C#, GDExtension (C++)"],
        ["Platform Desteği", "Masaüstü, Mobil, Web, Konsol"],
        ["Editör", "Yerleşik — kod, animasyon, shader, debugger dahil"],
        ["Yapı Taşları", "Node, Scene, Scene Tree, Signal"],
        ["Topluluk", "Aktif, büyüyen ve açık kaynak odaklı"],
    ]
    add_table_slide(slide, headers, rows, top=Inches(2.0))

    # ============================================================
    # SLIDE 33: SIRADAKİ ADIM / KAPANIŞ
    # ============================================================
    print("Slayt 33: Kapanış")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "Giriş Bölümü Tamamlandı!", font_size=48, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_divider(slide, Inches(2.8), Inches(4.5), Inches(4))

    items = [
        "Bu bölümde öğrendiklerimiz:",
        "",
        "•  Godot'nun ne olduğu ve neden tercih edildiği",
        "•  GDScript ile kod yazmaya başlamak",
        "•  Node, Scene, Scene Tree ve Signal kavramları",
        "•  Editör arayüzü ve beş ana ekran",
        "•  Tasarım felsefesi ve öğrenme kaynakları",
    ]
    add_bullet_slide_content(slide, items, left=Inches(2), top=Inches(3.2),
                              width=Inches(9), font_size=22, color=BODY_COLOR)

    add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
                 "Sıradaki Bölüm: Godot'yu kurup ilk projemizi oluşturuyoruz!",
                 font_size=24, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
                 "Görüşmek üzere...",
                 font_size=20, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # ---- SAVE ----
    prs.save(OUTPUT_FILE)
    print(f"\n✅ Sunum başarıyla oluşturuldu: {OUTPUT_FILE}")
    print(f"   Toplam slayt sayısı: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
