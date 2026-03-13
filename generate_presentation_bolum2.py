"""
Godot Engine Eğitim Serisi - Bölüm 2 (Giriş Kısmı): Sunum Oluşturucu
Node ve Sahne kavramlarının pekiştirilmesi - İlk Sahne öncesi teori.
Bölüm 1 ile aynı koyu tema.
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os
import io

# --- CONFIGURATION ---
BASE_DIR = r"e:\Data\Blogs_YouTube\Godot\Blog\fsadak.github.io"
ASSETS_DIR = os.path.join(BASE_DIR, "assets", "images")
OUTPUT_FILE = os.path.join(BASE_DIR, "Godot_Bolum2_Giris_Sunum.pptx")

# Colors - same dark theme as Bölüm 1
BG_COLOR = RGBColor(0x1E, 0x1E, 0x2E)
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
SUBTITLE_COLOR = RGBColor(0xB0, 0xB0, 0xCC)
BODY_COLOR = RGBColor(0xE0, 0xE0, 0xE8)
ACCENT_COLOR = RGBColor(0x47, 0x8C, 0xBF)
ACCENT2_COLOR = RGBColor(0x6C, 0x9E, 0x3A)
CODE_BG_COLOR = RGBColor(0x11, 0x11, 0x1B)
HIGHLIGHT_COLOR = RGBColor(0xFF, 0xCC, 0x33)
DIVIDER_COLOR = RGBColor(0x47, 0x8C, 0xBF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---- HELPERS (same as Bölüm 1) ----

def load_image_as_png(image_name):
    path = os.path.join(ASSETS_DIR, image_name)
    if not os.path.exists(path):
        print(f"  [UYARI] Görsel bulunamadı: {path}")
        return None
    try:
        img = Image.open(path)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf
    except Exception as e:
        print(f"  [HATA] Görsel yüklenemedi {image_name}: {e}")
        return None


def set_slide_bg(slide, color=BG_COLOR):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_divider(slide, top, left=Inches(0.8), width=Inches(2.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DIVIDER_COLOR
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BODY_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left=Inches(0.8), top=Inches(2.2),
                              width=Inches(11), font_size=20, color=BODY_COLOR,
                              spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if "**" in item:
            parts = item.split("**")
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.name = "Segoe UI"
                if j % 2 == 1:
                    run.font.bold = True
                    run.font.color.rgb = ACCENT_COLOR
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Segoe UI"

        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
    return txBox


def add_section_title(slide, title, subtitle=""):
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
                 title, font_size=36, color=TITLE_COLOR, bold=True)
    add_divider(slide, Inches(1.25))
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
                     subtitle, font_size=18, color=SUBTITLE_COLOR)


def add_image_to_slide(slide, image_name, left, top, width=None, height=None):
    buf = load_image_as_png(image_name)
    if buf is None:
        return None
    try:
        kwargs = {"image_file": buf, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        if not width and not height:
            kwargs["height"] = Inches(3.5)
        return slide.shapes.add_picture(**kwargs)
    except Exception as e:
        print(f"  [HATA] Görsel eklenemedi {image_name}: {e}")
        return None


def add_table_slide(slide, headers, rows, left=Inches(0.8), top=Inches(2.2)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    width = Inches(11.5)
    height = Inches(0.5) * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = TITLE_COLOR
            paragraph.font.bold = True
            paragraph.font.name = "Segoe UI"
            paragraph.alignment = PP_ALIGN.LEFT

    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = value
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x3C)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x33)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = BODY_COLOR
                paragraph.font.name = "Segoe UI"
                paragraph.alignment = PP_ALIGN.LEFT

    return table_shape


# ---- MAIN PRESENTATION BUILDER ----

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ============================================================
    # SLIDE 1: BAŞLIK
    # ============================================================
    print("Slayt 1: Başlık")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Godot Engine Eğitim Serisi", font_size=48, color=TITLE_COLOR, bold=True)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
                 "Bölüm 2: Node'lar ve Sahneler",
                 font_size=28, color=ACCENT_COLOR, bold=True)

    add_divider(slide, Inches(4.0), Inches(1), Inches(4))

    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 "İlk sahnemizi oluşturmadan önce temel yapı taşlarını pekiştirelim",
                 font_size=20, color=SUBTITLE_COLOR)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                 "Bu bölümde: Node nedir?  •  Sahne nedir?  •  Node ağaçları  •  Sahne yapısı",
                 font_size=16, color=RGBColor(0x88, 0x88, 0xAA))

    # ============================================================
    # SLIDE 2: İÇİNDEKİLER
    # ============================================================
    print("Slayt 2: İçindekiler")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "İçindekiler")

    items = [
        "1.  Hatırlayalım: Bölüm 1'den temel kavramlar",
        "2.  Node Nedir? — Detaylı bakış",
        "3.  Node'ların ortak özellikleri",
        "4.  Node ağacı ile karakter oluşturma",
        "5.  Sahne Nedir? — Detaylı bakış",
        "6.  Sahne özellikleri ve instancing",
        "7.  Godot editörü = Sahne editörü",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), font_size=22)

    # ============================================================
    # SLIDE 3: BÖLÜM 1'DEN HATIRLAYALIM
    # ============================================================
    print("Slayt 3: Bölüm 1'den Hatırlayalım")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Bölüm 1'den Hatırlayalım",
                      "Temel kavramlara hızlı bir bakış")

    headers = ["Kavram", "Tanım"]
    rows = [
        ["Node (Düğüm)", "Oyunun en küçük yapı taşı — her node'un belirli bir görevi var"],
        ["Sahne (Scene)", "Node'ların ağaç yapısında bir araya gelmesiyle oluşan birim"],
        ["Sahne Ağacı (Scene Tree)", "Tüm sahnelerin oluşturduğu bütünsel yapı"],
        ["Sinyal (Signal)", "Node'ların olay bazlı iletişim sistemi"],
    ]
    add_table_slide(slide, headers, rows, top=Inches(2.2))

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.6),
                 "Şimdi Node ve Sahne kavramlarını daha derinlemesine inceleyelim...",
                 font_size=18, color=ACCENT_COLOR)

    # ============================================================
    # SLIDE 4: BÖLÜM AYIRICI — NODE
    # ============================================================
    print("Slayt 4: Bölüm Ayırıcı — Node")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Node Nedir?", font_size=52, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Oyunun temel yapı taşları — tıpkı bir tarifin malzemeleri gibi",
                 font_size=24, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 5: NODE — TANIM
    # ============================================================
    print("Slayt 5: Node Tanım")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Node Nedir?",
                      "Oyunun temel yapı taşları")

    items = [
        "•  Node'lar oyunun **temel yapı taşlarıdır**",
        "•  Onlarca farklı türde node vardır:",
        "   ekranda görüntü çizen, ses çalayan, kamerayı temsil eden...",
        "",
        "•  Ekranda **görüntü çizmek** → Sprite2D, Label, AnimatedSprite2D",
        "•  **Ses çalmak** → AudioStreamPlayer",
        "•  **Kamera** → Camera2D, Camera3D",
        "•  **Fizik** → CharacterBody2D, RigidBody2D",
        "•  **Çarpışma** → CollisionShape2D",
        "•  Ve çok daha fazlası...",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6.5))

    add_image_to_slide(slide, "nodes_and_scenes_nodes.webp",
                       Inches(8), Inches(2.0), height=Inches(4.5))

    # ============================================================
    # SLIDE 6: NODE'LARIN ORTAK ÖZELLİKLERİ
    # ============================================================
    print("Slayt 6: Node'ların Ortak Özellikleri")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Tüm Node'ların Ortak Özellikleri")

    # 5 boxes for 5 properties
    props = [
        ("İsim", "Her node'un kendine\nözgü bir ismi vardır"),
        ("Özellikler", "Düzenlenebilir properties\n(Inspector'dan erişilir)"),
        ("Güncelleme", "Her kare (frame) için\ncallback alır"),
        ("Genişletme", "Yeni fonksiyonlarla\ngenişletilebilir"),
        ("Hiyerarşi", "Başka node'ların\nalt öğesi olabilir"),
    ]

    for i, (title, desc) in enumerate(props):
        x = Inches(0.5) + Inches(2.5) * i
        # Card background
        add_shape_rect(slide, x, Inches(2.2), Inches(2.3), Inches(3.2),
                       RGBColor(0x28, 0x28, 0x3C))
        # Number circle
        add_shape_rect(slide, x + Inches(0.1), Inches(2.3), Inches(0.4), Inches(0.4),
                       ACCENT_COLOR)
        add_text_box(slide, x + Inches(0.1), Inches(2.3), Inches(0.4), Inches(0.4),
                     str(i + 1), font_size=16, color=TITLE_COLOR, bold=True,
                     alignment=PP_ALIGN.CENTER)
        # Title
        add_text_box(slide, x + Inches(0.15), Inches(2.85), Inches(2.0), Inches(0.5),
                     title, font_size=18, color=ACCENT_COLOR, bold=True)
        # Description
        add_text_box(slide, x + Inches(0.15), Inches(3.4), Inches(2.0), Inches(1.5),
                     desc, font_size=14, color=BODY_COLOR)

    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.6),
                 "Son özellik çok önemli: Node'lar bir araya gelerek ağaç (tree) yapısı oluşturur!",
                 font_size=18, color=HIGHLIGHT_COLOR, bold=True)

    # ============================================================
    # SLIDE 7: NODE AĞACI İLE KARAKTER OLUŞTURMA
    # ============================================================
    print("Slayt 7: Node Ağacı ile Karakter")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Node Ağacı ile Karakter Oluşturma",
                      "Farklı görevlere sahip node'ları birleştirerek karmaşık davranışlar elde edebilirsiniz")

    items = [
        "Bir oyuncu karakteri şu node'lardan oluşabilir:",
        "",
        "•  **CharacterBody2D** — fizik ve hareket",
        "•  **Sprite2D** — karakterin görseli",
        "•  **Camera2D** — karakteri takip eden kamera",
        "•  **CollisionShape2D** — çarpışma alanı",
        "",
        "Her node'un **tek bir sorumluluğu** var.",
        "Birlikte çalışarak tam bir karakter oluşturuyorlar.",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6))

    add_image_to_slide(slide, "nodes_and_scenes_character_nodes.webp",
                       Inches(7.5), Inches(2.0), height=Inches(4.5))

    # ============================================================
    # SLIDE 8: BÖLÜM AYIRICI — SAHNE
    # ============================================================
    print("Slayt 8: Bölüm Ayırıcı — Sahne")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Sahne Nedir?", font_size=52, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Node ağaçlarından oluşan, kaydedilebilir ve tekrar kullanılabilir birimler",
                 font_size=24, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 9: SAHNE — TANIM
    # ============================================================
    print("Slayt 9: Sahne Tanım")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Sahne (Scene) Nedir?",
                      "Node'ları ağaç yapısında düzenleyince bir sahne elde edersiniz")

    items = [
        "•  Node'ları bir ağaç yapısında düzenlediğinizde buna **sahne** denir",
        "•  Kaydedildiğinde editörde **yeni bir node türü gibi** davranır",
        "•  Başka bir node'un alt öğesi olarak ekleyebilirsiniz",
        "•  Editörde iç yapısı gizlenir — **tek bir node olarak görünür**",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), width=Inches(6.5), font_size=20)

    add_image_to_slide(slide, "nodes_and_scenes_3d_scene_example.webp",
                       Inches(7.5), Inches(2.0), height=Inches(4.5))

    # ============================================================
    # SLIDE 10: SAHNE ÖZELLİKLERİ
    # ============================================================
    print("Slayt 10: Sahne Özellikleri")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Sahne Özellikleri")

    # Three cards
    cards = [
        ("Kök Node", "Her sahnenin her zaman\ntek bir kök node'u\n(root node) vardır",
         "🌳"),
        ("Kaydedilebilir", "Yerel diske kaydedilir\nve sonradan\nyüklenebilir (.tscn)",
         "💾"),
        ("Instancing", "İstediğin kadar örnek\n(instance) oluştur —\n10 farklı karakter yarat!",
         "✨"),
    ]

    for i, (title, desc, _icon) in enumerate(cards):
        x = Inches(0.8) + Inches(4) * i
        # Card background
        add_shape_rect(slide, x, Inches(2.2), Inches(3.6), Inches(3.5),
                       RGBColor(0x28, 0x28, 0x3C))
        # Accent bar on top of card
        add_shape_rect(slide, x, Inches(2.2), Inches(3.6), Pt(4), ACCENT_COLOR)
        # Title
        add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(3.0), Inches(0.6),
                     title, font_size=24, color=ACCENT_COLOR, bold=True)
        # Description
        add_text_box(slide, x + Inches(0.3), Inches(3.3), Inches(3.0), Inches(2.0),
                     desc, font_size=17, color=BODY_COLOR)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(12), Inches(0.6),
                 "Sahneler iç içe geçerek büyük ve karmaşık oyun dünyaları oluşturur",
                 font_size=18, color=SUBTITLE_COLOR)

    # ============================================================
    # SLIDE 11: GODOT EDİTÖRÜ = SAHNE EDİTÖRÜ
    # ============================================================
    print("Slayt 11: Godot Editörü = Sahne Editörü")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Godot Editörü = Sahne Editörü",
                      "2D, 3D ve kullanıcı arayüzü için zengin araçlar")

    items = [
        "•  Godot editörü özünde bir **sahne editörüdür**",
        "•  2D, 3D ve UI için zengin görsel araçlar sunar",
        "",
        "•  Bir Godot projesi istediğin kadar sahne içerebilir",
        "•  Ancak motorun başlatmak için **tek bir ana sahneye** ihtiyacı var",
        "",
        "•  Sahneler **.tscn** uzantısıyla kaydedilir",
        "•  Ana sahne **project.godot** dosyasında tanımlanır",
    ]
    add_bullet_slide_content(slide, items, top=Inches(2.0), font_size=20)

    # ============================================================
    # SLIDE 12: NODE VS SAHNE KARŞILAŞTIRMA
    # ============================================================
    print("Slayt 12: Node vs Sahne Karşılaştırma")
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Node ve Sahne: Karşılaştırma")

    headers = ["", "Node", "Sahne (Scene)"]
    rows = [
        ["Tanım", "En küçük yapı taşı", "Node'lardan oluşan ağaç yapısı"],
        ["Tek başına?", "Evet, tek başına çalışır", "En az bir kök node gerektirir"],
        ["Kaydetme", "Sahne içinde kaydedilir", ".tscn dosyası olarak kaydedilir"],
        ["Tekrar kullanım", "Her sahnede farklı ayarlarla", "Instance olarak çoğaltılır"],
        ["Örnek", "Sprite2D, Camera2D, Label", "Karakter, Silah, Seviye, Menü"],
    ]
    add_table_slide(slide, headers, rows, top=Inches(2.0))

    # ============================================================
    # SLIDE 13: KAPANIŞ
    # ============================================================
    print("Slayt 13: Kapanış")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, RGBColor(0x14, 0x14, 0x22))

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "Teoriyi Pekiştirdik!", font_size=48, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_divider(slide, Inches(2.8), Inches(4.5), Inches(4))

    items = [
        "Bu bölümde öğrendiklerimiz:",
        "",
        "•  Node'lar oyunun **en küçük yapı taşlarıdır**",
        "•  Her node'un ismi, özellikleri ve callback'leri vardır",
        "•  Node'lar **ağaç yapısı** oluşturarak sahneyi meydana getirir",
        "•  Sahneler **kaydedilebilir**, **iç içe geçebilir** ve **çoğaltılabilir**",
        "•  Godot editörü özünde bir **sahne editörüdür**",
    ]
    add_bullet_slide_content(slide, items, left=Inches(2), top=Inches(3.2),
                              width=Inches(9), font_size=22)

    add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
                 "Sıradaki: İlk sahnemizi oluşturuyoruz — Hello World!",
                 font_size=24, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
                 "Hadi başlayalım...",
                 font_size=20, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # ---- SAVE ----
    prs.save(OUTPUT_FILE)
    print(f"\n✅ Sunum başarıyla oluşturuldu: {OUTPUT_FILE}")
    print(f"   Toplam slayt sayısı: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
